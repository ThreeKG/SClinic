from __future__ import annotations

import sys
from pathlib import Path

# Avoid the local Pillow build when importing python-pptx.
sys.path = [p for p in sys.path if '/venv/' not in p]

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path('/Users/loay/Desktop/3kg/SClinic/Docs & Planing/deliverables')
OUT.mkdir(parents=True, exist_ok=True)

# Palette
DARK = '132238'
NAVY = '1B2B49'
TEAL = '1B8C8B'
GOLD = 'D7A643'
LIGHT = 'F5F7FA'
SLATE = '4A5668'
MUTED = '718096'
BORDER = 'D8DEE9'
WHITE = 'FFFFFF'
RED = 'B24A4A'

metrics = {
    'com': {'perf': 0.71, 'a11y': 0.61, 'bp': 0.54, 'seo': 0.83, 'fcp': '2.2s', 'lcp': '2.2s', 'tti': '9.9s', 'cls': '0.434', 'js': '311 KiB', 'cache': '2,172 KiB'},
    'ae': {'perf': 0.65, 'a11y': 0.81, 'bp': 0.54, 'seo': 0.83, 'fcp': '2.5s', 'lcp': '3.5s', 'tti': '12.1s', 'cls': '0.344', 'js': '429 KiB', 'cache': '1,936 KiB'},
}

roadmap = [
    ('Discover', 'Inventory content, routes, redirects'),
    ('Design', 'Define templates and design system'),
    ('Build', 'Implement CMS and pages'),
    ('QA + Launch', 'Validate performance and forms'),
    ('Stabilize', 'Monitor search, errors, conversions'),
]


def rgb(value: str):
    return RGBColor.from_string(value)


def set_bg(slide, fill):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(fill)


def textbox(slide, text, x, y, w, h, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font='Calibri', margin=0.04, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def card(slide, x, y, w, h, title, body, accent=TEAL, title_color=NAVY, body_color=SLATE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(WHITE)
    shape.line.color.rgb = rgb(BORDER)
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.07))
    strip.fill.solid(); strip.fill.fore_color.rgb = rgb(accent)
    strip.line.fill.background()
    textbox(slide, title, x+0.12, y+0.12, w-0.24, 0.25, size=14, bold=True, color=title_color)
    textbox(slide, body, x+0.12, y+0.42, w-0.24, h-0.5, size=10.2, color=body_color)


def metric(slide, x, y, w, h, label, value, note='', accent=TEAL, fill='FFFFFF'):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(BORDER)
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.06))
    strip.fill.solid(); strip.fill.fore_color.rgb = rgb(accent)
    strip.line.fill.background()
    textbox(slide, label.upper(), x+0.08, y+0.08, w-0.16, 0.16, size=9.2, bold=True, color=MUTED)
    textbox(slide, value, x+0.08, y+0.3, w-0.16, 0.26, size=19, bold=True, color=DARK)
    if note:
        textbox(slide, note, x+0.08, y+0.62, w-0.16, h-0.65, size=8.3, color=MUTED)


def bullets(slide, items, x, y, w, h, size=12.5, color=SLATE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.name = 'Calibri'
            r.font.size = Pt(size)
            r.font.color.rgb = rgb(color)
    return box


def footer(slide, page):
    textbox(slide, 'SClinic client deck', 0.35, 7.0, 4.8, 0.15, size=8.2, color=MUTED)
    textbox(slide, f'{page}', 12.45, 7.0, 0.4, 0.15, size=8.2, color=MUTED, align=PP_ALIGN.RIGHT)


def build(path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK)
    band = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.11))
    band.fill.solid(); band.fill.fore_color.rgb = rgb(GOLD); band.line.fill.background()
    textbox(s, 'SClinic Technical Audit', 0.7, 0.85, 6.4, 0.45, size=30, bold=True, color=WHITE)
    textbox(s, 'Client deck | evidence-backed recommendation to rebuild', 0.72, 1.35, 6.2, 0.24, size=14, color='DCE6F3')
    textbox(s, 'This deck summarizes the audit findings, business implications, and the recommended next step for the Turkish and Dubai sites.', 0.72, 1.68, 6.7, 0.55, size=16, color='F0F5FB')
    metric(s, 0.72, 4.65, 2.55, 1.0, 'Performance', '0.71 / 0.65', 'Lighthouse scores', accent=TEAL, fill='17304C')
    metric(s, 3.42, 4.65, 2.55, 1.0, 'Risk', 'High', 'Legacy stack + JS bug', accent=RED, fill='17304C')
    metric(s, 6.12, 4.65, 2.55, 1.0, 'Decision', 'Rebuild', 'Recommended direction', accent=GOLD, fill='17304C')
    card(s, 8.85, 0.95, 3.75, 1.45, 'Evidence basis', 'Headers, robots/sitemap captures, browser inspection, console output, and Lighthouse JSON for both sites.', accent=GOLD, title_color=NAVY)
    card(s, 8.85, 2.55, 3.75, 1.45, 'Core issues', 'Legacy PHP 7.4.33 stack, runtime JS exception, repeated hero content, and weak public hardening.', accent=TEAL, title_color=NAVY)
    card(s, 8.85, 4.15, 3.75, 1.45, 'Outcome', 'Proceed with a phased rebuild that improves speed, trust, SEO clarity, and maintainability.', accent=RED, title_color=NAVY)

    # 2 summary
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT)
    textbox(s, 'Executive summary', 0.55, 0.3, 4, 0.35, size=22, bold=True)
    textbox(s, 'Both sites are functional, but they share a legacy technical foundation and a homepage pattern that is too heavy, too repetitive, and too fragile. The current state is good enough to operate, but not ideal for a premium healthcare brand.', 0.55, 0.8, 6.1, 1.0, size=16, color=SLATE)
    card(s, 0.55, 2.0, 5.8, 1.1, 'Decision', 'Rebuild the platform rather than layering more fixes onto the current stack.', accent=GOLD, title_color=NAVY)
    bullets(s, [
        'Legacy stack: nginx + PHP/7.4.33 + PleskLin',
        'Same JavaScript exception on both homepages',
        'The .ae sitemap still contains Turkish legacy URLs',
        'Lighthouse reports layout shifts, unused JS, and cache waste',
    ], 0.72, 3.45, 5.4, 2.2)
    card(s, 6.75, 0.8, 5.9, 2.1, 'What this means', 'Incremental patching will keep the current problems alive: slower interaction, more maintenance burden, weaker accessibility, and less clarity for users choosing a consultation path.', accent=TEAL, title_color=NAVY)
    card(s, 6.75, 3.2, 5.9, 1.7, 'Commercial implication', 'A cleaner rebuild gives the team a better platform for conversions, localization, and SEO growth.', accent=RED, title_color=NAVY)
    card(s, 6.75, 5.15, 5.9, 1.2, 'Recommendation', 'Approve a phased rebuild with strict migration rules and performance budgets.', accent=GOLD, title_color=NAVY)

    # 3 evidence dashboard
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    textbox(s, 'Evidence dashboard', 0.55, 0.25, 4.8, 0.35, size=22, bold=True)
    chart_data = CategoryChartData()
    chart_data.categories = ['Performance', 'Accessibility', 'Best Practices', 'SEO']
    chart_data.add_series('sclinic.com.tr', [metrics['com']['perf'], metrics['com']['a11y'], metrics['com']['bp'], metrics['com']['seo']])
    chart_data.add_series('sclinic.ae', [metrics['ae']['perf'], metrics['ae']['a11y'], metrics['ae']['bp'], metrics['ae']['seo']])
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(0.85), Inches(7.2), Inches(4.45), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.maximum_scale = 1.0
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.major_unit = 0.2
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.plots[0].has_data_labels = True
    card(s, 7.95, 0.95, 4.7, 1.0, 'Runtime bug', 'Both homepages log the same JavaScript exception in aliaygir.js.', accent=RED, title_color=NAVY)
    card(s, 7.95, 2.1, 4.7, 1.0, 'Cache waste', 'Lighthouse estimates 2.2 MiB and 1.9 MiB of cache savings.', accent=GOLD, title_color=NAVY)
    card(s, 7.95, 3.25, 4.7, 1.0, 'Layout shifts', 'Each page triggers 15 layout shifts and visible instability.', accent=TEAL, title_color=NAVY)
    card(s, 7.95, 4.4, 4.7, 1.45, 'Stack proof', 'Headers expose nginx, PHP/7.4.33, and PleskLin. Browser inspection shows legacy jQuery-era assets plus GTM, Google Ads, and Facebook Pixel.', accent=GOLD, title_color=NAVY)
    metric(s, 0.75, 5.55, 2.45, 0.9, 'com.tr unused JS', '311 KiB', 'Est. savings', accent=TEAL)
    metric(s, 3.35, 5.55, 2.45, 0.9, 'ae unused JS', '429 KiB', 'Est. savings', accent=TEAL)
    metric(s, 5.95, 5.55, 2.45, 0.9, 'com.tr CLS', '0.434', 'Too high', accent=RED)
    metric(s, 8.55, 5.55, 2.45, 0.9, 'ae CLS', '0.344', 'Still high', accent=RED)

    # 4 performance
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT)
    textbox(s, 'Performance', 0.55, 0.25, 3, 0.35, size=22, bold=True)
    card(s, 0.55, 0.85, 5.9, 5.6, 'What the numbers say', 'First paint is acceptable, but interactivity is slow, layout stability is weak, and the sites are carrying too much unused script and image weight.', accent=TEAL, title_color=NAVY)
    table = [
        ['Metric', 'sclinic.com.tr', 'sclinic.ae'],
        ['Performance', '0.71', '0.65'],
        ['FCP', '2.2s', '2.5s'],
        ['LCP', '2.2s', '3.5s'],
        ['TTI', '9.9s', '12.1s'],
        ['CLS', '0.434', '0.344'],
    ]
    x, y = 6.75, 0.95
    w, h = 5.95, 2.45
    t = s.shapes.add_table(len(table), 3, Inches(x), Inches(y), Inches(w), Inches(h)).table
    widths = [1.9, 1.0, 1.0]
    for i, ww in enumerate(widths):
        t.columns[i].width = Inches(ww)
    for r, row in enumerate(table):
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb(NAVY if r == 0 else ('F7FAFC' if r % 2 else WHITE))
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = 'Calibri'
                    run.font.size = Pt(10)
                    run.font.bold = r == 0
                    run.font.color.rgb = rgb(WHITE if r == 0 else DARK)
    card(s, 6.75, 3.7, 5.95, 1.2, 'Operational impact', 'More bytes, more scripts, more shifts, and slower interaction readiness create friction for mobile users and weaken conversion confidence.', accent=GOLD, title_color=NAVY)
    metric(s, 6.75, 5.15, 2.75, 1.0, 'com.tr cache savings', '2.2 MiB', 'Estimated', accent=TEAL)
    metric(s, 9.95, 5.15, 2.75, 1.0, 'ae cache savings', '1.9 MiB', 'Estimated', accent=TEAL)

    # 5 security / stack
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    textbox(s, 'Security and stack posture', 0.55, 0.25, 5.5, 0.35, size=22, bold=True)
    card(s, 0.55, 0.9, 3.9, 2.0, 'Public posture', 'The captured response does not show a modern hardening baseline such as CSP, HSTS, X-Frame-Options, Referrer-Policy, or Permissions-Policy.', accent=RED, title_color=NAVY)
    card(s, 0.55, 3.15, 3.9, 2.0, 'Cookie / script risk', 'The session cookie is exposed without visible HttpOnly in the capture, and the page relies on multiple third-party tags that add privacy and performance cost.', accent=GOLD, title_color=NAVY)
    card(s, 0.55, 5.4, 3.9, 1.0, 'Runtime bug', 'A null addEventListener exception appears in aliaygir.js.', accent=TEAL, title_color=NAVY)
    card(s, 4.8, 0.9, 3.9, 2.0, 'Legacy foundation', 'The stack is still anchored in nginx + PHP/7.4.33 + PleskLin, with old jQuery, carousel, lightbox, and lazyload assets.', accent=TEAL, title_color=NAVY)
    card(s, 4.8, 3.15, 3.9, 2.0, 'Recommended now', 'If the current stack stays live, add the security baseline, remove the JS exception, and cut unused scripts immediately.', accent=GOLD, title_color=NAVY)
    card(s, 4.8, 5.4, 3.9, 1.0, 'Risk level', 'Medium-high for hardening; high for maintainability.', accent=RED, title_color=NAVY)
    card(s, 9.05, 0.9, 3.75, 5.5, 'Executive takeaway', 'This is not evidence of active compromise. It is evidence of a dated public posture that should be treated as a rebuild trigger, not a long-term operating model.', accent=GOLD, title_color=NAVY)

    # 6 SEO / localization
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT)
    textbox(s, 'SEO and localization', 0.55, 0.25, 5.5, 0.35, size=22, bold=True)
    card(s, 0.55, 0.9, 4.1, 5.75, 'What is working', 'robots.txt allows crawling, both sites expose sitemaps, and SEO scores are acceptable. The foundation exists, but it is not yet cleanly organized for two markets.', accent=TEAL, title_color=NAVY)
    metric(s, 5.0, 1.0, 2.0, 0.95, '.com.tr URLs', '60', 'Unique sitemap URLs', accent=TEAL)
    metric(s, 7.2, 1.0, 2.0, 0.95, '.ae URLs', '62', 'Unique sitemap URLs', accent=TEAL)
    metric(s, 9.4, 1.0, 2.0, 0.95, 'Locale drift', 'Yes', 'Turkish routes remain', accent=RED)
    card(s, 5.0, 2.35, 6.2, 1.25, 'Localization issue', 'The .ae sitemap still includes Turkish legacy URLs such as s-clinic-eryaman-ankaraman, s-clinic-cayyolu, and s-clinic-ankara. That suggests reuse rather than market-specific architecture.', accent=GOLD, title_color=NAVY)
    card(s, 5.0, 3.85, 6.2, 1.25, 'SEO implication', 'Mixed-market routing can dilute relevance, confuse users, and make search intent harder to read.', accent=RED, title_color=NAVY)
    card(s, 5.0, 5.35, 6.2, 1.25, 'Rebuild objective', 'Separate locale content cleanly and rebuild metadata, structured data, and URL conventions from the ground up.', accent=TEAL, title_color=NAVY)
    card(s, 11.35, 2.35, 1.45, 4.25, 'Takeaway', 'Use market-specific content models, canonical rules, and sitemap governance.', accent=GOLD, title_color=NAVY)

    # 7 UX / accessibility
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    textbox(s, 'UX and accessibility', 0.55, 0.25, 5.5, 0.35, size=22, bold=True)
    card(s, 0.55, 0.9, 5.2, 5.75, 'User experience evidence', 'The homepage uses repeated hero slides, multiple similar CTAs, and a dense content stack. The result is visual noise rather than a clear path to consultation.', accent=RED, title_color=NAVY)
    bullets(s, [
        'Repeated hero content weakens message clarity.',
        'The primary CTA is not dominant enough at first glance.',
        'Content blocks feel stacked rather than prioritized.',
        'Mobile usability is affected by layout shifts and tap-target issues.',
    ], 0.8, 2.15, 4.6, 2.0)
    card(s, 6.0, 0.9, 6.7, 2.05, 'Accessibility proof', 'Lighthouse flags missing alt text, missing main landmarks, links without discernible text, insufficient contrast, and touch targets that are too small.', accent=GOLD, title_color=NAVY)
    card(s, 6.0, 3.15, 6.7, 1.45, 'Conversion impact', 'When the hierarchy is unclear, users work harder to find consultation, location, and trust information — especially on mobile.', accent=TEAL, title_color=NAVY)
    card(s, 6.0, 4.85, 6.7, 1.8, 'Rebuild target', 'A cleaner information architecture, stronger CTA hierarchy, and accessible components should be core requirements.', accent=RED, title_color=NAVY)

    # 8 roadmap
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT)
    textbox(s, 'Rebuild recommendation and roadmap', 0.55, 0.25, 7.2, 0.35, size=22, bold=True)
    card(s, 0.55, 0.9, 5.65, 5.8, 'Proposed solution', 'Build a modern healthcare platform with a clean content model, reusable components, SEO-safe templates, performance budgets, and a localized structure for Turkish and English markets. The goal is a platform that is faster, safer, and easier to operate.', accent=TEAL, title_color=NAVY)
    card(s, 6.35, 0.9, 6.15, 1.2, 'Core stack', 'Modern CMS or custom WordPress block architecture, CDN-backed hosting, responsive images, analytics governance, and secure defaults.', accent=GOLD, title_color=NAVY)
    x = 6.35
    y = 2.35
    widths = [1.15, 1.15, 1.15, 1.15, 1.15]
    colors = [TEAL, GOLD, NAVY, TEAL, RED]
    for i, (phase, desc) in enumerate(roadmap):
        card(s, x + i * 1.22, y, 1.12, 3.0, phase, desc, accent=colors[i], title_color=NAVY)
    card(s, 6.35, 5.65, 6.15, 0.95, 'Success criteria', 'Faster pages, clearer consultation funnels, cleaner localization, and a platform the team can maintain confidently.', accent=RED, title_color=NAVY)

    prs.save(str(path))


if __name__ == '__main__':
    build(OUT / 'SClinic_Client_Deck.pptx')
